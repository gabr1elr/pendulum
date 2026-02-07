import numpy as np
import matplotlib.pyplot as plt
import imageio
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

output_dir = Path('pendul_simplu_output')
output_dir.mkdir(exist_ok=True)

# Parametri fizici
g = 9.81        # acceleratia gravitationala (m/s^2)
L = 1.0         # lungimea firului (m)
theta0 = 0.3    # unghi initial (radiani)
omega0 = 0.0    # viteza unghiulara initiala (rad/s)


t_max = 10
dt = 0.01
t = np.arange(0, t_max, dt)

# Ecuatia miscarii
def pendulum_eq(theta, omega, dt):
    omega_new = omega - (g/l) * np.sin(theta) * dt
    theta_new = theta + omega_new * dt
    return theta_new, omega_new


theta = np.zeros_like(t)
omega = np.zeros_like(t)
theta[0] = theta0
omega[0] = omega0

for i in range(1, len(t)):
    theta[i], omega[i] = pendulum_eq(theta[i-1], omega[i-1], dt)

x = L * np.sin(theta)
y = -L * np.cos(theta)

# Animatie pendul
frames = []
for i in range(0, len(t), 10):
    fig, ax = plt.subplots(figsize=(3,3))
    ax.plot([0, x[i]], [0, y[i]], 'k-', lw=2)
    ax.plot(x[i], y[i], 'ro', markersize=10)
    ax.set_xlim(-L-0.2, L+0.2)
    ax.set_ylim(-L-0.2, 0.2)
    ax.set_aspect('equal')
    ax.axis('off')
    ax.set_title(f"Pendul simplu (t={t[i]:.1f}s)")
    fname = output_dir / f"pend_{i:04d}.png"
    fig.savefig(fname, bbox_inches='tight')
    plt.close(fig)
    frames.append(imageio.v2.imread(fname))

gif_path = output_dir / "pendul.gif"
imageio.mimsave(gif_path, frames, fps=20)

#Grafic
fig, ax = plt.subplots()
ax.plot(t, theta)
ax.set_xlabel("Timp (s)")
ax.set_ylabel("Unghi θ (rad)")
ax.set_title("Evoluția unghiului pendulului în timp")
ax.grid(True)
theta_fig = output_dir / "theta_vs_time.png"
fig.savefig(theta_fig, bbox_inches='tight')
plt.close(fig)

# Prezentare pptx
prs = Presentation()

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Oscilațiile unui pendul simplu"
slide.placeholders[1].text = "Model fizic, simulare numerică și animație — realizat în Python"

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "1. Teorie de bază"
body = slide.placeholders[1].text_frame
body.text = "Ecuația mișcării pendulului:\n d²θ/dt² + (g/l)·sinθ = 0"
p = body.add_paragraph()
p.text = "Pentru unghiuri mici: sinθ ≈ θ → mișcare armonică simplă."
p.level = 1
p = body.add_paragraph()
p.text = "Soluția aproximativă: θ(t) = θ₀·cos(sqrt(g/l)·t)"#sqrt fiind radical
p.level = 1

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "2. Animația pendulului"
slide.shapes.add_picture(str(gif_path), Inches(1.5), Inches(1.5), width=Inches(6))

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "3. Graficul θ(t)"
slide.shapes.add_picture(str(theta_fig), Inches(1), Inches(1.5), width=Inches(8))

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "4. Concluzii"
body = slide.placeholders[1].text_frame
body.text = "Pendulul simplu este un exemplu de mișcare oscilatorie periodică."
p = body.add_paragraph()
p.text = "Pentru unghiuri mici, perioada este independentă de amplitudine."
p.level = 1
p = body.add_paragraph()
p.text = "Pentru unghiuri mari, apare abaterea de la armonicitate."
p.level = 1

pptx_path = output_dir / "Prezentare_Pendul_Simplu.pptx"
prs.save(pptx_path)

print(f"Fisiere generate in: {output_dir.resolve()}")

